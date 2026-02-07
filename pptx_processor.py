"""
PPTX Processor - OPTIMIZED
--------------------------
Fast processing with batch translation per slide.
"""

import logging
from typing import List, Tuple
from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.shapes.base import BaseShape
from pptx.table import Table
from pptx.text.text import TextFrame

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PPTXProcessor:
    def __init__(self, translator, status_callback=None):
        self.translator = translator
        self.status_callback = status_callback
        self.stats = {
            'slides_processed': 0,
            'shapes_processed': 0,
            'text_runs_translated': 0,
            'tables_processed': 0,
            'notes_translated': 0,
            'groups_traversed': 0,
            'errors': []
        }
    
    def _update_status(self, message):
        if self.status_callback:
            self.status_callback(message)
        logger.info(message)
    
    def process_file(self, input_path: str, output_path: str) -> dict:
        try:
            self._update_status("Loading presentation...")
            prs = Presentation(input_path)
            num_slides = len(prs.slides)
            self._update_status(f"Loaded {num_slides} slides - using fast batch mode")
            
            for slide_idx, slide in enumerate(prs.slides):
                try:
                    self._update_status(f"Slide {slide_idx + 1}/{num_slides}...")
                    self._process_slide_batch(slide)
                    self.stats['slides_processed'] += 1
                except Exception as e:
                    error_msg = f"Error on slide {slide_idx + 1}: {e}"
                    logger.error(error_msg)
                    self.stats['errors'].append(error_msg)
            
            self._update_status("Saving translated presentation...")
            prs.save(output_path)
            
            self.stats['translator_stats'] = self.translator.get_stats()
            return self.stats
            
        except Exception as e:
            logger.error(f"Error processing file: {e}")
            self.stats['errors'].append(str(e))
            raise
    
    def _process_slide_batch(self, slide) -> None:
        """
        OPTIMIZED: Collect all texts from slide, batch translate, then apply.
        Much faster than translating one text at a time.
        """
        # Collect all text locations
        text_locations: List[Tuple] = []  # (paragraph, runs, original_texts, combined_text)
        
        # Gather from all shapes
        for shape in slide.shapes:
            self._collect_texts(shape, text_locations)
        
        # Gather from notes
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                for para in slide.notes_slide.notes_text_frame.paragraphs:
                    runs = list(para.runs)
                    if runs:
                        orig = [r.text for r in runs]
                        combined = ''.join(orig)
                        if combined.strip():
                            text_locations.append((para, runs, orig, combined))
        except:
            pass
        
        if not text_locations:
            return
        
        # Batch translate all texts at once
        all_texts = [loc[3] for loc in text_locations]
        translated_texts = self.translator.translate_batch(all_texts)
        
        # Apply translations back
        for (para, runs, orig_texts, orig_combined), translated in zip(text_locations, translated_texts):
            if translated and translated != orig_combined:
                self._redistribute_text_to_runs(runs, orig_texts, translated)
                self.stats['text_runs_translated'] += len(runs)
    
    def _collect_texts(self, shape: BaseShape, text_locations: List) -> None:
        """Recursively collect all translatable text from a shape."""
        try:
            if isinstance(shape, GroupShape):
                self.stats['groups_traversed'] += 1
                for child in shape.shapes:
                    self._collect_texts(child, text_locations)
                return
            
            if shape.has_table:
                self._collect_table_texts(shape.table, text_locations)
                self.stats['tables_processed'] += 1
                return
            
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    runs = list(para.runs)
                    if runs:
                        orig = [r.text for r in runs]
                        combined = ''.join(orig)
                        if combined.strip() and len(combined.strip()) >= 2:
                            text_locations.append((para, runs, orig, combined))
            
            self.stats['shapes_processed'] += 1
            
        except Exception as e:
            self.stats['errors'].append(f"Shape error: {e}")
    
    def _collect_table_texts(self, table: Table, text_locations: List) -> None:
        """Collect all text from table cells."""
        try:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        for para in cell.text_frame.paragraphs:
                            runs = list(para.runs)
                            if runs:
                                orig = [r.text for r in runs]
                                combined = ''.join(orig)
                                if combined.strip() and len(combined.strip()) >= 2:
                                    text_locations.append((para, runs, orig, combined))
        except Exception as e:
            self.stats['errors'].append(f"Table error: {e}")
    
    def _redistribute_text_to_runs(self, runs, original_texts, translated_text) -> None:
        """Redistribute translated text back to runs proportionally."""
        total_len = sum(len(t) for t in original_texts)
        
        if total_len == 0:
            if runs:
                runs[0].text = translated_text
            return
        
        trans_len = len(translated_text)
        pos = 0
        
        for i, (run, orig) in enumerate(zip(runs, original_texts)):
            if i == len(runs) - 1:
                run.text = translated_text[pos:]
            else:
                prop = len(orig) / total_len
                chars = int(trans_len * prop)
                end = pos + chars
                
                # Break at word boundary if possible
                if end < trans_len:
                    space = translated_text.rfind(' ', pos, end + 10)
                    if space > pos:
                        end = space + 1
                
                run.text = translated_text[pos:end]
                pos = end
